"""Adapter: PowerPoint .pptx -> list[Element].  ⚠️ FALLBACK ONLY (deprecated as primary).

The project's COMMITTED docx/pptx path is MinerU's native office backend (scripts/parse_office.py
-> content_list -> from_mineru). On pptx, fair measurement showed MinerU ties or slightly trails
Tika and is comparable to this adapter; MinerU was chosen for its chunk-ready content_list schema
and pipeline unification with PDF/scanned. Kept only as a zero-dependency fallback.

Slides ARE pages -> page-grouped mode (one slide = one chunk). Title placeholder -> heading;
body/tables/pictures/charts follow. Use Chunker(page_grouped={"pptx"}).chunk(..., doc_type="pptx")."""
from __future__ import annotations

from ..types import Element


def _is_title(shape):
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER
        if not shape.is_placeholder:
            return False
        return shape.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
    except Exception:
        return False


def _table_html(tbl):
    rows = []
    for r in tbl.rows:
        cells = "".join(f"<td>{(c.text or '').strip()}</td>" for c in r.cells)
        rows.append(f"<tr>{cells}</tr>")
    return "<table>" + "".join(rows) + "</table>"


def _leaf_shapes(shapes):
    """Yield non-group shapes, recursing into GROUP — slide.shapes does NOT recurse, so text
    nested in grouped shapes is otherwise silently dropped (review F2)."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for s in shapes:
        try:
            if s.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from _leaf_shapes(s.shapes)
                continue
        except Exception:
            pass
        yield s


def _ordered_shapes(slide):
    """Title first, then remaining leaf shapes (groups flattened) in document order."""
    leaves = list(_leaf_shapes(slide.shapes))
    titles = [s for s in leaves if _is_title(s)]
    rest = [s for s in leaves if not _is_title(s)]
    return titles + rest


def from_pptx(path) -> list[Element]:
    """Read a .pptx into Element[] (page = slide index, so page-grouped mode yields one chunk
    per slide). Speaker notes are intentionally excluded."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(path))
    out, idx = [], 0
    for pno, slide in enumerate(prs.slides):
        for shape in _ordered_shapes(slide):
            try:
                if shape.has_table:
                    out.append(Element(idx=idx, kind="table", text=None,
                                       table_body=_table_html(shape.table), page=pno))
                    idx += 1
                    continue
                if getattr(shape, "has_chart", False):
                    out.append(Element(idx=idx, kind="chart", text=None, page=pno))
                    idx += 1
                    continue
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    out.append(Element(idx=idx, kind="image", text=None, page=pno))
                    idx += 1
                    continue
                if shape.has_text_frame:
                    title = _is_title(shape)
                    for para in shape.text_frame.paragraphs:
                        txt = (para.text or "").strip()   # para.text renders <a:br/> line breaks (review F4)
                        if not txt:
                            continue
                        out.append(Element(idx=idx, kind="text", text=txt,
                                           text_level=1 if title else None, page=pno))
                        idx += 1
            except Exception:
                continue
    return out
